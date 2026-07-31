"""Export PowerPoint — remplit le template `cr_vt_template.pptx` avec les
données d'une session. Les slides 1-7, 9, 10 ne sont pas modifiées.

Slides remplies :
- 8  : Tableau Commandes (split en 2 tables, 24×6 et 25×6)
- 11 : Tableau date global (5 × N nuits)
- 12 : Phasage EEG/Rails complet (N+2 × 8) — étend dynamiquement
- 13-16 : Phasage par semaine S1..S4 (2 tables chacune)
- 17 : Tableau date caméras (4 × N nuits cam)
- 18 : Phasage caméras (N+2 × 5)
- 19 : Détail caméras par allée (split en 2 tables, étend dynamiquement)
- 20 : Phasage full consolidé (N+3 × 10)

Style : couleurs par position de la nuit dans la semaine (cf. night_color_hex).
"""
from __future__ import annotations
import copy
import re
import copy
from datetime import datetime
from pathlib import Path
from io import BytesIO

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import logging

logger = logging.getLogger(__name__)

TEMPLATE_PATH = Path(__file__).parent / "templates" / "cr_vt_template.pptx"
LOGISTIQUE_SLIDE_PATH = Path(__file__).parent / "templates" / "slide_logistique.pptx"

# Marqueur de version pour debug deploy — incrémenter à chaque changement majeur.
# Visible dans le header HTTP `X-PPTX-Version` de la réponse d'export.
__PPTX_VERSION__ = "2026-07-31-v27-logistique+nightcolors"

# Palette par position dans la semaine (alignée Excel)
# (iter48j) Cycle de 4 couleurs distinctes appliqué au n° absolu de la nuit
# (bleu → jaune → rose → vert → bleu → …). Avant : DBEAFE était en position 0
# ET position 3, ce qui produisait des paires bleu-bleu au changement de semaine.
WEEK_COLORS_HEX = ["#DBEAFE", "#FEF3C7", "#FECACA", "#DCFCE7"]
WHITE = "#FFFFFF"
HEADER_BG = "#1F2937"
SUBHEADER_BG = "#F3F4F6"


def _pos_in_week(nuit: int, weeks: list | None) -> int:
    if not nuit:
        return 0
    if not weeks:
        return int(nuit)
    remaining = int(nuit)
    for w in weeks:
        ww = int(w or 0)
        if remaining <= ww:
            return remaining
        remaining -= ww
    return remaining


def _color_for_night(n: int, weeks: list | None) -> str:
    # (iter48k) La couleur reflète la POSITION dans la semaine (règle métier
    # utilisateur) : 1ère nuit d'une semaine = bleu, 2e = jaune, 3e = rose,
    # 4e = vert. Une semaine peut n'avoir que 2 ou 3 nuits — dans ce cas la
    # semaine suivante recommence à bleu. La palette a été corrigée
    # (position 3 était bleu au lieu de vert, ce qui produisait bleu-bleu).
    if not n:
        return WHITE
    pos = _pos_in_week(n, weeks)
    if not pos:
        return WHITE
    return WEEK_COLORS_HEX[(pos - 1) % len(WEEK_COLORS_HEX)]


def _set_cell_text(cell, value, *, bold=False, italic=False, align="center", size=None, color=None, fill_rgb=None):
    """Remplace le contenu d'une cellule en préservant approximativement le
    style. On garde le premier paragraphe + run existants si présents.

    fill_rgb : tuple (r, g, b) optionnel pour la couleur de fond de la cellule.
    """
    if fill_rgb is not None:
        try:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(*fill_rgb)
        except Exception:
            pass
    tf = cell.text_frame
    # Conserve les paragraphes/runs existants pour garder police/style si on
    # peut, sinon on rebuilde.
    if tf.paragraphs:
        p = tf.paragraphs[0]
        # Nettoie tous les autres paragraphes
        for extra_p in tf.paragraphs[1:]:
            extra_p._p.getparent().remove(extra_p._p)
        # Garde le 1er run si possible
        if p.runs:
            r = p.runs[0]
            # Supprime runs additionnels
            for extra_r in p.runs[1:]:
                extra_r._r.getparent().remove(extra_r._r)
            r.text = str(value) if value is not None else ""
        else:
            r = p.add_run()
            r.text = str(value) if value is not None else ""
        if bold:
            r.font.bold = True
        if italic:
            r.font.italic = True
        if size:
            r.font.size = Pt(size)
        if color:
            r.font.color.rgb = RGBColor.from_string(color.lstrip("#"))
        if align == "center":
            p.alignment = PP_ALIGN.CENTER
        elif align == "left":
            p.alignment = PP_ALIGN.LEFT
        elif align == "right":
            p.alignment = PP_ALIGN.RIGHT


def _force_cell_text_color(cell, hex_color: str = "000000"):
    """Force la couleur du texte sur TOUS les runs d'une cellule en
    supprimant tout `<a:solidFill>` hérité du template (via clonage de
    lignes) puis en réinjectant un srgbClr propre au niveau `<a:rPr>`.

    Nécessaire car `_set_cell_text` (via python-pptx) ne remplace pas
    toujours fiablement la couleur héritée du template quand celui-ci
    contient plusieurs runs ou des couleurs scheme/thème."""
    NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
    hx = hex_color.lstrip("#").upper()
    tc = cell._tc
    for r in tc.iter(f"{{{NS}}}r"):
        rpr = r.find(f"{{{NS}}}rPr")
        if rpr is None:
            rpr = etree.SubElement(r, f"{{{NS}}}rPr")
            # rPr doit être 1er enfant du run
            r.insert(0, rpr)
        # Supprime tout solidFill existant
        for child in list(rpr):
            if etree.QName(child).localname == "solidFill":
                rpr.remove(child)
        # Injecte un solidFill/srgbClr propre en tête du rPr
        sf = etree.SubElement(rpr, f"{{{NS}}}solidFill")
        srgb = etree.SubElement(sf, f"{{{NS}}}srgbClr")
        srgb.set("val", hx)


def _set_cell_fill(cell, hex_color: str):
    """Définit la couleur de fond d'une cellule via XML direct (python-pptx
    n'expose pas table_cell.fill correctement pour les TableCells)."""
    if not hex_color or hex_color == WHITE:
        return
    tc = cell._tc
    tcPr = tc.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}tcPr')
    if tcPr is None:
        tcPr = etree.SubElement(tc, '{http://schemas.openxmlformats.org/drawingml/2006/main}tcPr')
    # Supprime les fills existants
    for child in list(tcPr):
        tag = etree.QName(child).localname
        if tag in ("solidFill", "noFill", "gradFill", "blipFill", "pattFill"):
            tcPr.remove(child)
    nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    fill = etree.SubElement(tcPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
    srgb = etree.SubElement(fill, '{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
    srgb.set("val", hex_color.lstrip("#").upper())
    _ = nsmap


def _clone_last_row(table):
    """Duplique la dernière ligne du tableau et l'insère à la fin.
    Permet d'étendre dynamiquement une table sans casser le style."""
    tbl = table._tbl
    rows = tbl.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}tr')
    if not rows:
        return None
    last = rows[-1]
    new = copy.deepcopy(last)
    # Vide le texte des cellules de la nouvelle ligne
    for tc in new.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}tc'):
        for p in tc.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}txBody/{http://schemas.openxmlformats.org/drawingml/2006/main}p'):
            for r in p.findall('{http://schemas.openxmlformats.org/drawingml/2006/main}r'):
                t = r.find('{http://schemas.openxmlformats.org/drawingml/2006/main}t')
                if t is not None:
                    t.text = ""
    last.addnext(new)


def _ensure_table_size(table, n_rows: int):
    """Étend une table jusqu'à n_rows en clonant la dernière ligne."""
    while len(table.rows) < n_rows:
        _clone_last_row(table)


def _unmerge_row_cells(table, row_idx: int, n_cols: int):
    """Retire tous les attributs de fusion (gridSpan/hMerge/vMerge/rowSpan)
    sur les `n_cols` premières cellules d'une ligne. À utiliser AVANT
    d'écrire dans des cellules d'un tableau étendu, pour s'assurer qu'aucune
    cellule ne reste invisible (fantôme) du fait de merges hérités du
    template."""
    NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    tbl = table._tbl
    trs = tbl.findall(f'{{{NS}}}tr')
    if row_idx >= len(trs):
        return
    tcs = trs[row_idx].findall(f'{{{NS}}}tc')
    for tc in tcs[:n_cols]:
        tc.attrib.pop('gridSpan', None)
        tc.attrib.pop('hMerge', None)
        tc.attrib.pop('rowSpan', None)
        tc.attrib.pop('vMerge', None)


def _clone_last_col(table):
    """Duplique la dernière colonne du tableau (gridCol + tc dans chaque tr).
    Conserve les styles de la dernière colonne — mais nettoie les attributs
    de fusion (gridSpan / hMerge) hérités du template : sinon les nouvelles
    cellules seraient fusionnées avec la voisine et resteraient invisibles
    (bug observé sur le slide 11 lorsque le template a > 16 gridCol figés)."""
    NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    tbl = table._tbl
    grid = tbl.find(f'{{{NS}}}tblGrid')
    if grid is None:
        return
    grid_cols = grid.findall(f'{{{NS}}}gridCol')
    if not grid_cols:
        return
    new_gc = copy.deepcopy(grid_cols[-1])
    grid.append(new_gc)
    for tr in tbl.findall(f'{{{NS}}}tr'):
        tcs = tr.findall(f'{{{NS}}}tc')
        if not tcs:
            continue
        new_tc = copy.deepcopy(tcs[-1])
        # (v24) Retire les attributs de fusion hérités — sinon la cellule
        # clonée serait considérée comme continuation de la précédente et ne
        # s'afficherait pas (fantôme).
        new_tc.attrib.pop('gridSpan', None)
        new_tc.attrib.pop('rowSpan', None)
        new_tc.attrib.pop('hMerge', None)
        new_tc.attrib.pop('vMerge', None)
        for p in new_tc.findall(f'.//{{{NS}}}txBody/{{{NS}}}p'):
            for r in p.findall(f'{{{NS}}}r'):
                t = r.find(f'{{{NS}}}t')
                if t is not None:
                    t.text = ""
        tr.append(new_tc)


def _ensure_table_cols(table, n_cols: int, label_cols: int = 1):
    """Étend une table jusqu'à n_cols en clonant la dernière colonne.
    Les `label_cols` premières conservent leur largeur ; les autres
    (data_cols) se partagent équitablement la somme des largeurs des
    anciennes data_cols pour rester dans la largeur totale d'origine."""
    NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    tbl = table._tbl
    grid = tbl.find(f'{{{NS}}}tblGrid')
    if grid is None:
        return
    grid_cols = grid.findall(f'{{{NS}}}gridCol')
    orig_n = len(grid_cols)
    if orig_n >= n_cols:
        return
    try:
        orig_data_total = sum(int(gc.get('w')) for gc in grid_cols[label_cols:])
    except (TypeError, ValueError):
        orig_data_total = None
    while len(table.columns) < n_cols:
        _clone_last_col(table)
    if orig_data_total is not None:
        new_grid_cols = grid.findall(f'{{{NS}}}gridCol')
        data_cols = new_grid_cols[label_cols:]
        if data_cols:
            new_w = orig_data_total // len(data_cols)
            for gc in data_cols:
                gc.set('w', str(new_w))


def _trim_table_cols(table, n_cols: int):
    """Supprime les colonnes au-delà de n_cols (gridCol + tc dans chaque tr).
    Sert à retirer les colonnes vides résiduelles du template (à droite)."""
    NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    tbl = table._tbl
    grid = tbl.find(f'{{{NS}}}tblGrid')
    if grid is None:
        return
    grid_cols = grid.findall(f'{{{NS}}}gridCol')
    if len(grid_cols) <= n_cols:
        return
    for gc in grid_cols[n_cols:]:
        grid.remove(gc)
    for tr in tbl.findall(f'{{{NS}}}tr'):
        tcs = tr.findall(f'{{{NS}}}tc')
        for tc in tcs[n_cols:]:
            tr.remove(tc)


def _merge_title_row(table, n_cols: int, row_idx: int = 0):
    """Force la ligne `row_idx` à être une unique cellule fusionnée sur les
    n_cols premières colonnes (bandeau de titre). Corrige le cas où le template
    a un gridSpan figé (< n_cols) après ajout de colonnes → cellules vides qui
    dépassent à droite."""
    NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    tbl = table._tbl
    trs = tbl.findall(f'{{{NS}}}tr')
    if row_idx >= len(trs):
        return
    tcs = trs[row_idx].findall(f'{{{NS}}}tc')
    if len(tcs) < n_cols:
        return
    tcs[0].set('gridSpan', str(n_cols))
    tcs[0].attrib.pop('hMerge', None)
    for tc in tcs[1:n_cols]:
        tc.set('hMerge', '1')
        tc.attrib.pop('gridSpan', None)


def _merge_cells_range(table, row_idx: int, start_col: int, end_col: int):
    """Fusionne les cellules [start_col..end_col] inclus sur la ligne row_idx.
    Utilisé pour créer des sous-bandeaux (ex : slide 20 « Phasage étiquettes »
    couvre cols 0-7, « Nuit » 8-9, « Phasage caméras » 10-12)."""
    NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    tbl = table._tbl
    trs = tbl.findall(f'{{{NS}}}tr')
    if row_idx >= len(trs):
        return
    tcs = trs[row_idx].findall(f'{{{NS}}}tc')
    if end_col >= len(tcs) or start_col > end_col:
        return
    span = end_col - start_col + 1
    tcs[start_col].set('gridSpan', str(span))
    tcs[start_col].attrib.pop('hMerge', None)
    for tc in tcs[start_col + 1:end_col + 1]:
        tc.set('hMerge', '1')
        tc.attrib.pop('gridSpan', None)



def _remove_table_row(table, row_idx: int):
    """Supprime la ligne d'index row_idx (utile pour retirer un en-tête/bandeau)."""
    trs = table._tbl.findall(qn('a:tr'))
    if 0 <= row_idx < len(trs):
        trs[row_idx].getparent().remove(trs[row_idx])


def _fmt_date(iso: str | None) -> str:
    if not iso:
        return ""
    try:
        return datetime.strptime(iso, "%Y-%m-%d").strftime("%d/%m/%Y")
    except Exception:
        return iso


def _get_tables(slide):
    return [sh for sh in slide.shapes if sh.has_table]


def _set_col_widths_by_ratio(table, ratios: list[int]):
    """Répartit la largeur totale actuelle de la table selon `ratios`
    (un poids par colonne). Préserve la largeur totale d'origine."""
    try:
        total = sum(int(c.width) for c in table.columns)
    except (TypeError, ValueError):
        return
    s = sum(ratios) or 1
    for ci, w in enumerate(ratios):
        if ci < len(table.columns):
            table.columns[ci].width = Emu(int(total * w / s))


def _compact_layout(prs):
    """Post-traitement mise en page (gain de place) :
      - remonte le titre des slides qui contiennent un tableau,
      - réduit la police de la phrase d'avertissement / mention de bas de page
        sur toutes les slides où elle apparaît."""
    FOOT_KEYS = (
        "allées numérotées", "correspondantes sont à retrouver",
        "proprietary and confidential", "signalétique du magasin",
    )
    for slide in prs.slides:
        has_table = any(getattr(sh, "has_table", False) for sh in slide.shapes)
        try:
            title = slide.shapes.title
        except Exception:
            title = None
        for sh in slide.shapes:
            if not getattr(sh, "has_text_frame", False):
                continue
            low = (sh.text_frame.text or "").lower()
            if any(k in low for k in FOOT_KEYS):
                for p in sh.text_frame.paragraphs:
                    for r in p.runs:
                        r.font.size = Pt(7)
            if has_table and title is not None and sh == title:
                try:
                    sh.top = Emu(90000)
                except Exception:
                    pass


def _place_table(shape, left, top, width, row_h):
    """Repositionne/dimensionne le GraphicFrame d'une table pour compacter la
    mise en page (mesures reprises de la maquette de référence utilisateur) :
      - place le tableau (left, top) et fixe sa largeur totale `width`,
      - redistribue les colonnes proportionnellement pour remplir `width`,
      - applique une hauteur de ligne uniforme `row_h` (compacte).
    La hauteur totale s'adapte au nombre réel de lignes (donc au nb de nuits)."""
    try:
        shape.left = Emu(int(left))
        shape.top = Emu(int(top))
        shape.width = Emu(int(width))
        t = shape.table
        cols = list(t.columns)
        cur = [int(c.width) for c in cols]
        s = sum(cur) or 1
        for c, w in zip(cols, cur):
            c.width = Emu(int(width * w / s))
        trs = t._tbl.findall(qn('a:tr'))
        for tr in trs:
            tr.set('h', str(int(row_h)))
        shape.height = Emu(int(row_h * max(1, len(trs))))
    except Exception:
        pass


# En-têtes + ratios de largeur pour les tables "Phasage par nuit" (slide 12 +
# semaines) après éclatement de la colonne SA en SA 1.5 / 2.1 / frz (+ magasin)
# et retrait de la colonne Caméras (les caméras ont leurs propres slides).
def _phasage_headers(hide_sa_mag: bool) -> list[str]:
    base = ["Nuit", "Date", "Secteur/Rayon", "Allées", "EEG", "Rails ES",
            "SA 1.5", "SA 2.1", "SA 2.1 frz"]
    if not hide_sa_mag:
        base.append("SA magasin")
    return base


def _phasage_ratios(hide_sa_mag: bool) -> list[int]:
    base = [7, 9, 17, 20, 8, 8, 7, 7, 8]
    if not hide_sa_mag:
        base.append(8)
    return base




# ===================================================================
# Slide 8 — Commandes (split en 2 tables 24×6 et 25×6)
# ===================================================================
# Headers et largeurs pour les NOUVELLES tables 10-cols créées via add_table().
_RECAP_COL_HEADERS = [
    "Type", "Réf.", "Désignation", "Total", "Spare", "Flèche",
    "Signalétique", "Saisonnier", "Total", "Total + MOQ",
]
# Largeurs finalisées : Type 8% pour "Fixation" tienne sur une ligne partout.
_RECAP_COL_WEIGHTS = [8, 6, 28, 6, 6, 6, 11, 10, 7, 12]


def _set_recap_col_widths(table, total_emu: int):
    s = sum(_RECAP_COL_WEIGHTS)
    for ci, w in enumerate(_RECAP_COL_WEIGHTS):
        table.columns[ci].width = Emu(int(total_emu * w / s))


def _set_cell_margins_zero(cell):
    """Marges internes ultra-réduites pour maximiser l'espace texte ET
    minimiser la hauteur de ligne effective."""
    tcPr = cell._tc.get_or_add_tcPr()
    tcPr.set('marL', '36000')   # 0.04 inch L/R
    tcPr.set('marR', '36000')
    tcPr.set('marT', '0')       # 0 padding T/B → row height au plus près du texte
    tcPr.set('marB', '0')


def _fill_slide_8(slide, recap_rows: list):
    """27/02/2026 v9 — On SUPPRIME les 2 tables 6-cols du template et on CRÉE
    2 nouvelles tables 10-cols via slide.shapes.add_table(). Cela génère un
    XML 100 % compatible PowerPoint Desktop (comme un copier-coller Excel→PPT).
    """
    # 1) Localise les positions/tailles des 2 tables existantes pour les recréer
    old_tables = _get_tables(slide)
    if len(old_tables) < 2:
        return
    placements = []
    for gframe in old_tables[:2]:
        placements.append({
            "left": gframe.left, "top": gframe.top,
            "width": gframe.width, "height": gframe.height,
        })
    # 2) Supprime les anciennes tables
    for gframe in old_tables[:2]:
        sp = gframe._element
        sp.getparent().remove(sp)
    # 3) Filtre les lignes (exclut VCare et vides)
    rows = [r for r in recap_rows
            if r.get("kind") != "empty" and r.get("type") != "VCare"]
    # 4) Calcule la répartition rows → t1 / t2
    cap_per_table = 24
    n_t1 = min(len(rows), cap_per_table)
    rest = rows[n_t1:]
    # 5) Crée les 2 nouvelles tables — HAUTEUR = N rows × row_height_fixe
    # pour forcer des lignes très compactes (PowerPoint ignore le `h` sur tr
    # mais respecte la hauteur totale de la table).
    ROW_HEIGHT_EMU = 120000  # ≈ 0.13 inch — assez pour 7pt + padding 0
    n_rows_t1 = 1 + n_t1
    n_rows_t2 = max(1 + len(rest), 2)
    h_t1 = ROW_HEIGHT_EMU * n_rows_t1
    h_t2 = ROW_HEIGHT_EMU * n_rows_t2
    t1_shape = slide.shapes.add_table(n_rows_t1, 10,
                                       placements[0]["left"], placements[0]["top"],
                                       placements[0]["width"], h_t1)
    t2_shape = slide.shapes.add_table(n_rows_t2, 10,
                                       placements[1]["left"], placements[1]["top"],
                                       placements[1]["width"], h_t2)
    t1, t2 = t1_shape.table, t2_shape.table
    # 6) Largeurs proportionnelles
    _set_recap_col_widths(t1, placements[0]["width"])
    _set_recap_col_widths(t2, placements[1]["width"])
    # 7) Headers + data
    _write_recap_header(t1, 0)
    _write_recap_header(t2, 0)
    for i, r in enumerate(rows[:n_t1]):
        _write_recap_row(t1, i + 1, r)
    for i, r in enumerate(rest):
        _write_recap_row(t2, i + 1, r)
    # Hauteur de ligne FORCÉE — même valeur sur chaque tr ET sur la hauteur totale.
    for tbl in (t1, t2):
        for tr in tbl._tbl.findall(qn('a:tr')):
            tr.set('h', str(ROW_HEIGHT_EMU))


def _write_recap_header(table, row_idx: int):
    """Header gris/gras sur 10 cols."""
    for c, label in enumerate(_RECAP_COL_HEADERS):
        align = "left" if c < 3 else "right"
        cell = table.cell(row_idx, c)
        _set_cell_text(cell, label, bold=True, align=align, size=7,
                       fill_rgb=(0xE5, 0xE7, 0xEB))
        _set_cell_margins_zero(cell)


def _write_recap_row(table, row_idx, r):
    is_section = r.get("kind") == "section"
    if is_section:
        # Banner : on FUSIONNE les 10 cellules en une seule (comme un copier-coller
        # Excel), puis on remplit avec le nom de section + fond bleu clair.
        first = table.cell(row_idx, 0)
        last = table.cell(row_idx, 9)
        first.merge(last)
        _set_cell_text(first, (r.get("type") or ""),
                       bold=True, align="left", size=7,
                       fill_rgb=(0xDD, 0xEB, 0xF7))
        _set_cell_margins_zero(first)
        return
    vals = [
        (r.get("type", ""), "left", False),
        (r.get("reference", ""), "left", False),
        # Force la désignation "Inclineur" simple même si le dataset stocke
        # encore l'ancien texte long "Inclineur (1 par rail 1320/...)".
        ("Inclineur" if r.get("kind") == "inclineur" else r.get("designation", ""), "left", False),
        (_num(r.get("quantite")), "right", False),
        (_num(r.get("spare")), "right", False),
        (_num(r.get("fleche")), "right", False),
        (_num(r.get("signaletique")), "right", False),
        (_num(r.get("saisonnier")), "right", False),
        (_num(r.get("total_plus_spare")), "right", True),
        (("—" if r.get("total_moq") == "—" else _num(r.get("total_moq"))), "right", True),
    ]
    for c, (txt, align, bold) in enumerate(vals):
        cell = table.cell(row_idx, c)
        _set_cell_text(cell, txt, bold=bold, align=align, size=7)
        _set_cell_margins_zero(cell)


def _clear_row(table, row_idx):
    for c in range(len(table.columns)):
        _set_cell_text(table.cell(row_idx, c), "", size=8)


def _replace_nb_nuits_in_title(slide, nb_nuits: int):
    """Met à jour le titre d'une slide qui contient '(X nuits)' avec la
    bonne valeur de nb_nuits. Cherche le pattern '(<nombre> nuits)' dans
    tous les text-frames de la slide et le remplace tout en préservant
    le formatage (run par run).

    Ex: 'Tableau phasage EEG et rails par nuit (14 nuits)' →
        'Tableau phasage EEG et rails par nuit (12 nuits)' si nb_nuits=12.
    """
    if not nb_nuits or nb_nuits <= 0:
        return
    pattern = re.compile(r"\(\s*\d+\s+nuits?\s*\)")
    replacement = f"({nb_nuits} nuits)"
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        for para in sh.text_frame.paragraphs:
            # Concatène le texte du paragraphe pour détecter le pattern
            full = "".join(r.text or "" for r in para.runs)
            if not pattern.search(full):
                continue
            new_full = pattern.sub(replacement, full)
            if new_full == full:
                continue
            # On remet tout le texte sur le premier run en gardant son
            # formatage. Les runs suivants sont vidés.
            runs = list(para.runs)
            if runs:
                runs[0].text = new_full
                for r in runs[1:]:
                    r.text = ""


def _num(v):
    if v is None or v == "":
        return ""
    try:
        f = float(v)
        return f"{int(f)}" if f.is_integer() else f"{f:.2f}"
    except (ValueError, TypeError):
        return str(v)


# ===================================================================
# Slide 11 — Tableau date global (5 × N nuits)
# Lignes : [_, Date, EEG, Caméra, SA]   Colonnes : [label, Nuit 1, ..., Nuit N]
# ===================================================================
def _fill_slide_11(slide, totals_by_nuit, dates_map, weeks, all_nights: list[int]):
    tables = _get_tables(slide)
    if not tables:
        return
    _shape = tables[0]
    t = _shape.table
    # Étend/tronque le nb de colonnes pour couvrir TOUTES les nuits (plus de
    # troncature : les nuits 17-20 étaient perdues auparavant).
    needed_cols = 1 + len(all_nights)
    _ensure_table_cols(t, needed_cols, label_cols=1)
    _trim_table_cols(t, needed_cols)

    # (v24) Nettoie les merges hérités sur les 4 lignes du tableau avant
    # d'écrire — évite les cellules fantômes 17-18 fusionnées avec leurs
    # voisines (bug reproduit avec un phasage 18 nuits en prod).
    for r in range(4):
        _unmerge_row_cells(t, r, needed_cols)

    # Row 0 : header (Nuit X)
    _set_cell_text(t.cell(0, 0), "", bold=True, size=6)
    for i, n in enumerate(all_nights):
        _set_cell_text(t.cell(0, i + 1), f"Nuit {n}", bold=True, align="center", size=6)
        _set_cell_fill(t.cell(0, i + 1), _color_for_night(n, weeks))

    # Lignes : Date / EEG ES+SA / SA magasin (ligne « Caméra » retirée).
    labels = ["Date", "EEG ES+SA", "SA magasin"]
    _ensure_table_size(t, 1 + len(labels))
    for li, lab in enumerate(labels):
        r = li + 1
        italic = (lab == "SA magasin")
        _set_cell_text(t.cell(r, 0), lab, bold=True, align="left", size=6,
                       italic=italic, color=("#6B7280" if italic else "#000000"))
        for i, n in enumerate(all_nights):
            tot = totals_by_nuit.get(n, {})
            if lab == "Date":
                val = _fmt_date(dates_map.get(str(n)))
            elif lab == "EEG ES+SA":
                val = _num(tot.get("eeg", 0))
            else:
                val = _num(tot.get("sa", 0))
            _set_cell_text(t.cell(r, i + 1), val, size=6,
                           bold=(lab == "EEG ES+SA"),
                           italic=italic,
                           color=("#6B7280" if italic else "#000000"),
                           align="center")
            _set_cell_fill(t.cell(r, i + 1), _color_for_night(n, weeks))
            # Force couleur (le template a des runs violets/gris hérités).
            if italic:
                _force_cell_text_color(t.cell(r, i + 1), "#6B7280")
            else:
                _force_cell_text_color(t.cell(r, i + 1), "#000000")
    # Force couleur sur la ligne d'en-tête aussi.
    _force_cell_text_color(t.cell(0, 0), "#000000")
    for i in range(len(all_nights)):
        _force_cell_text_color(t.cell(0, i + 1), "#000000")
    # Supprime les lignes résiduelles (ancienne ligne Caméra / SA en trop)
    while len(t.rows) > 1 + len(labels):
        last_tr = t._tbl.findall(qn('a:tr'))[-1]
        last_tr.getparent().remove(last_tr)
    # Position/dimensions compactes (bandeau fin en haut — maquette de référence)
    _place_table(_shape, 1822174, 574934, 8347544, 161670)


# ===================================================================
# Slide 12 — Phasage EEG/Rails complet (N+2 × 8) — étend
# ===================================================================
def _fill_slide_12(slide, nuit_es_data, weeks, hide_sa_mag=False):
    tables = _get_tables(slide)
    if not tables:
        return
    _shape = tables[0]
    t = _shape.table
    nights_sorted = sorted(nuit_es_data.keys())
    n_nights = len(nights_sorted)
    needed_rows = 2 + n_nights + 1  # bandeau + header + nuits + TOTAL
    _ensure_table_size(t, needed_rows)
    # Colonnes cible : SA détaillées par type (SA 1.5, SA 2.1, SA 2.1 frz,
    # 4.2/4.2 WP) ; colonne Caméras ajoutée. Pas de SA magasin.
    headers = ["Nuit", "Date", "Secteur/Rayon", "Allées", "EEG ES", "Rails ES", "SA 1.5", "SA 2.1", "SA 2.1 frz", "4.2/4.2 WP", "Caméras"]
    ncols = len(headers)
    # Date élargie pour tenir "27/07/2026" sur UNE seule ligne à 9pt.
    ratios = [8, 11, 16, 20, 7, 7, 6, 6, 7, 6, 6]
    _ensure_table_cols(t, ncols, label_cols=1)
    _trim_table_cols(t, ncols)
    _set_col_widths_by_ratio(t, ratios)
    _merge_title_row(t, ncols, row_idx=0)
    _set_cell_text(t.cell(0, 0), "Récap par nuit", bold=True, align="center", size=11)
    for ci, h in enumerate(headers):
        _set_cell_text(t.cell(1, ci), h, bold=True, align="center", size=9)
    # Force le texte noir sur la ligne d'en-tête (le template a des runs
    # violet/gris hérités qui subsistent après clone_col).
    for ci in range(ncols):
        _force_cell_text_color(t.cell(1, ci), "#000000")
    # Cumuls pour la ligne TOTAL
    tot = {"eeg": 0, "rails_es": 0, "sa_inst_15": 0, "sa_inst_21": 0,
           "sa_inst_freezer": 0, "sa_inst_42": 0, "cam": 0}
    allees_planifiees = set()
    for i, n in enumerate(nights_sorted):
        r = i + 2
        d = nuit_es_data[n]
        _set_cell_text(t.cell(r, 0), f"Nuit {n}", bold=True, size=9)
        _set_cell_text(t.cell(r, 1), _fmt_date(d.get("date")), size=9)
        _set_cell_text(t.cell(r, 2), d.get("sr", ""), align="left", size=8)
        _set_cell_text(t.cell(r, 3), d.get("allees_str", ""), align="left", size=8)
        _set_cell_text(t.cell(r, 4), _num(d.get("eeg", 0)), bold=True, size=9)
        _set_cell_text(t.cell(r, 5), _num(d.get("rails_es", 0)), size=9)
        _set_cell_text(t.cell(r, 6), _num(d.get("sa_inst_15", 0) or ""), size=9)
        _set_cell_text(t.cell(r, 7), _num(d.get("sa_inst_21", 0) or ""), size=9)
        _set_cell_text(t.cell(r, 8), _num(d.get("sa_inst_freezer", 0) or ""), size=9)
        _set_cell_text(t.cell(r, 9), _num(d.get("sa_inst_42", 0) or ""), size=9)
        _set_cell_text(t.cell(r, 10), _num(d.get("cam", 0) or ""), size=9)
        color = _color_for_night(n, weeks)
        for ci in range(ncols):
            _set_cell_fill(t.cell(r, ci), color)
            # Force le texte en noir (le template a des couleurs
            # violet/gris héritées sur les colonnes SA/Caméras).
            _force_cell_text_color(t.cell(r, ci), "#000000")
        # Cumul TOTAL
        for k in tot.keys():
            v = d.get(k if k != "eeg" else "eeg", 0)
            try:
                tot[k] += int(v or 0)
            except (ValueError, TypeError):
                pass
        for a in str(d.get("allees_str") or "").split(","):
            a = a.strip()
            if a:
                allees_planifiees.add(a)
    # Ligne TOTAL (cf. modèle de référence utilisateur du 10/07/2026)
    tr = 2 + n_nights
    _set_cell_text(t.cell(tr, 0), "TOTAL", bold=True, align="center", size=10)
    _set_cell_text(t.cell(tr, 1), "", size=9)
    _set_cell_text(t.cell(tr, 2), "", size=9)
    _set_cell_text(t.cell(tr, 3), f"{len(allees_planifiees)} allées planifiées", bold=True, align="center", size=9)
    _set_cell_text(t.cell(tr, 4), _num(tot["eeg"]), bold=True, size=10)
    _set_cell_text(t.cell(tr, 5), _num(tot["rails_es"]), bold=True, size=10)
    _set_cell_text(t.cell(tr, 6), _num(tot["sa_inst_15"] or ""), bold=True, size=10)
    _set_cell_text(t.cell(tr, 7), _num(tot["sa_inst_21"] or ""), bold=True, size=10)
    _set_cell_text(t.cell(tr, 8), _num(tot["sa_inst_freezer"] or ""), bold=True, size=10)
    _set_cell_text(t.cell(tr, 9), _num(tot["sa_inst_42"] or ""), bold=True, size=10)
    _set_cell_text(t.cell(tr, 10), _num(tot["cam"] or ""), bold=True, size=10)
    for ci in range(ncols):
        _set_cell_fill(t.cell(tr, ci), "FEF3C7")  # jaune pâle (comme slide 19)
        _force_cell_text_color(t.cell(tr, ci), "#000000")
    # Supprime les lignes vides finales (au-delà de needed_rows)
    while len(t.rows) > needed_rows:
        last_tr = t._tbl.findall(qn('a:tr'))[-1]
        last_tr.getparent().remove(last_tr)
    # Position/dimensions compactes (maquette de référence)
    # row_h réduit pour faire tenir jusqu'à 20 nuits + TOTAL sans dépasser la slide.
    _place_table(_shape, 666206, 983106, 10115008, 230000)


# ===================================================================
# Slides 13-16 — Par semaine : tableau détaillé (Nuit/Date/SR/Allées/EEG/
# Rails ES + colonnes SA à installer dynamiques) + ligne « Sous-total S{n} ».
# ===================================================================
def _fill_slide_week(slide, week_index: int, week_nights: list[int],
                     nuit_es_data, totals_by_nuit, dates_map, weeks, hide_sa_mag=False):
    tables = _get_tables(slide)
    if not tables:
        return
    # On garde la table détaillée (plus de colonnes) et on supprime l'éventuelle
    # petite table transposée du template.
    shapes_sorted = sorted(tables, key=lambda sh: len(sh.table.columns), reverse=True)
    t_shape = shapes_sorted[0]
    t = t_shape.table
    for sh in shapes_sorted[1:]:
        sh._element.getparent().remove(sh._element)
    _remove_table_row(t, 0)  # retire l'éventuel bandeau/en-tête noir du template

    # Colonnes fixes : SA 1.5 / SA 2.1 / SA 2.1 frz / 4.2/4.2 WP + Caméras
    # (11 colonnes strictes cf. modèle de référence utilisateur du 10/07/2026 —
    # la colonne "SA magasin" n'apparaît PAS dans les semaines).
    sa_cols = [
        ("SA 1.5", "sa_inst_15", False),
        ("SA 2.1", "sa_inst_21", False),
        ("SA 2.1 frz", "sa_inst_freezer", False),
        ("4.2/4.2 WP", "sa_inst_42", False),
    ]
    include_cam = True

    headers = ["Nuit", "Date", "Secteur/Rayon", "Allées", "EEG ES", "Rails ES"] + [c[0] for c in sa_cols]
    if include_cam:
        headers.append("Caméras")
    ncols = len(headers)
    # Ratios (Date large pour dates complètes)
    ratios = [7, 12, 18, 20, 8, 8] + [7] * len(sa_cols)
    if include_cam:
        ratios.append(8)
    n_data = len(week_nights)
    needed = 1 + n_data + 1  # header + nuits + sous-total
    _ensure_table_size(t, needed)
    _ensure_table_cols(t, ncols, label_cols=1)
    _trim_table_cols(t, ncols)
    _set_col_widths_by_ratio(t, ratios)

    for ci, h in enumerate(headers):
        _set_cell_text(t.cell(0, ci), h, bold=True, align="center", size=9)

    sub = {"eeg": 0, "rails_es": 0, "cam": 0}
    for c in sa_cols:
        sub[c[1]] = 0
    for i, n in enumerate(week_nights):
        r = i + 1
        d = nuit_es_data.get(n, {})
        _set_cell_text(t.cell(r, 0), f"Nuit {n}", bold=True, size=9)
        _set_cell_text(t.cell(r, 1), _fmt_date(d.get("date") or dates_map.get(str(n))), size=9)
        _set_cell_text(t.cell(r, 2), d.get("sr", ""), align="left", size=8)
        _set_cell_text(t.cell(r, 3), d.get("allees_str", ""), align="left", size=8)
        eeg = int(round(d.get("eeg", 0) or 0)); rails = int(round(d.get("rails_es", 0) or 0))
        _set_cell_text(t.cell(r, 4), _num(eeg), bold=True, size=9)
        _set_cell_text(t.cell(r, 5), _num(rails), size=9)
        sub["eeg"] += eeg; sub["rails_es"] += rails
        for j, (h, key, ital) in enumerate(sa_cols):
            v = int(round(d.get(key, 0) or 0))
            sub[key] += v
            _set_cell_text(t.cell(r, 6 + j), _num(v or ""), size=9, italic=ital,
                           color=("#6B7280" if ital else None))
        if include_cam:
            cam_val = int(round((totals_by_nuit.get(n, {}) or {}).get("cam", 0) or 0))
            sub["cam"] += cam_val
            _set_cell_text(t.cell(r, 6 + len(sa_cols)), _num(cam_val or ""), size=9)
        color = _color_for_night(n, weeks)
        for ci in range(ncols):
            _set_cell_fill(t.cell(r, ci), color)
            # Force texte noir sur cellules non-magasin (magasin reste italique/gris).
            is_italic_mag = (6 <= ci < 6 + len(sa_cols) and sa_cols[ci - 6][2])
            if not is_italic_mag:
                _force_cell_text_color(t.cell(r, ci), "#000000")

    # Ligne « Sous-total S{n} »
    sr = 1 + n_data
    _set_cell_text(t.cell(sr, 0), f"Sous-total S{week_index}", bold=True, align="left", size=9)
    _set_cell_text(t.cell(sr, 1), "", size=9)
    _set_cell_text(t.cell(sr, 2), "", size=9)
    _set_cell_text(t.cell(sr, 3), "", size=9)
    _set_cell_text(t.cell(sr, 4), _num(sub["eeg"]), bold=True, size=9)
    _set_cell_text(t.cell(sr, 5), _num(sub["rails_es"]), bold=True, size=9)
    for j, (h, key, ital) in enumerate(sa_cols):
        _set_cell_text(t.cell(sr, 6 + j), _num(sub[key] or ""), bold=True, size=9, italic=ital,
                       color=("#6B7280" if ital else None))
    if include_cam:
        _set_cell_text(t.cell(sr, 6 + len(sa_cols)), _num(sub["cam"] or ""), bold=True, size=9)
    for ci in range(ncols):
        _set_cell_fill(t.cell(sr, ci), "F3F4F6")
        # Force texte noir sur sous-total (sauf colonne magasin italique).
        is_italic_mag = (ci >= 6 and (ci - 6) < len(sa_cols) and sa_cols[ci - 6][2])
        if not is_italic_mag:
            _force_cell_text_color(t.cell(sr, ci), "#000000")

    # Force aussi le texte noir sur la ligne d'en-tête (le template a du
    # texte blanc/violet hérité sur les colonnes SA après clone_col).
    for ci in range(ncols):
        _force_cell_text_color(t.cell(0, ci), "#000000")

    while len(t.rows) > needed:
        last_tr = t._tbl.findall(qn('a:tr'))[-1]
        last_tr.getparent().remove(last_tr)
    # Position/dimensions : tableau large (12 colonnes) sans écraser le
    # titre à gauche. On laisse ~3" pour le titre (2 lignes max) puis
    # tableau sur ~9.5" de largeur.
    TABLE_LEFT = 2900000
    _place_table(t_shape, TABLE_LEFT, 500000, 9100000, 175000)
    # Rétrécir titre pour ne pas chevaucher le tableau + réduire la police
    # pour rester sur 2 lignes.
    for sh in slide.shapes:
        if sh.has_text_frame and "Plan de phasage" in (sh.text_frame.text or ""):
            try:
                new_w = TABLE_LEFT - int(sh.left) - 80000
                if new_w > 914400:
                    sh.width = new_w
                    sh.text_frame.word_wrap = True
                # Réduit la taille de police du titre pour rester compact
                # (2 lignes idéalement, 3 max) sur toutes les slides semaines
                for p in sh.text_frame.paragraphs:
                    for r in p.runs:
                        if r.font.size and r.font.size.pt > 18:
                            r.font.size = Pt(18)
            except (TypeError, ValueError):
                pass
            break


# ===================================================================
# Slide 17 — Tableau date caméras (4 × N)
# ===================================================================
def _fill_slide_17(slide, totals_by_nuit, dates_map, cam_nights: list[int], weeks):
    tables = _get_tables(slide)
    if not tables:
        return
    _shape = tables[0]
    t = _shape.table
    # Étend dynamiquement le tableau pour accueillir toutes les nuits caméras
    # (1 colonne label + N colonnes nuits). Avant : tronqué à cur_cols-1 nuits.
    _ensure_table_cols(t, 1 + len(cam_nights))
    _trim_table_cols(t, 1 + len(cam_nights))
    cur_cols = len(t.columns)
    nights = cam_nights[: cur_cols - 1]
    _set_cell_text(t.cell(0, 0), "", size=11)
    for i, n in enumerate(nights):
        _set_cell_text(t.cell(0, i + 1), f"Nuit {n}", bold=True, align="center", size=12,
                       color="#000000")
        _set_cell_fill(t.cell(0, i + 1), _color_for_night(n, weeks))
        _force_cell_text_color(t.cell(0, i + 1), "#000000")
    # Lignes : Date / Caméra (ligne « EEG » retirée — réf. utilisateur).
    labels = ["Date", "Caméra"]
    _ensure_table_size(t, 1 + len(labels))
    for li, lab in enumerate(labels):
        r = li + 1
        _set_cell_text(t.cell(r, 0), lab, bold=True, align="left", size=11, color="#000000")
        for i, n in enumerate(nights):
            tot = totals_by_nuit.get(n, {})
            if lab == "Date":
                val = _fmt_date(dates_map.get(str(n)))
            else:
                val = _num(tot.get("cam", 0))
            _set_cell_text(t.cell(r, i + 1), val, align="center", size=11,
                           bold=(lab == "Caméra"), color="#000000")
            _set_cell_fill(t.cell(r, i + 1), _color_for_night(n, weeks))
            _force_cell_text_color(t.cell(r, i + 1), "#000000")
    # Supprime les lignes résiduelles du template (au-delà de 1 + labels)
    while len(t.rows) > 1 + len(labels):
        last_tr = t._tbl.findall(qn('a:tr'))[-1]
        last_tr.getparent().remove(last_tr)
    # Position/dimensions compactes (bandeau — maquette de référence).
    # top abaissé pour ne pas chevaucher le titre de la slide.
    _place_table(_shape, 1952045, 1250000, 7461250, 187000)


# ===================================================================
# Slide 18 — Phasage caméras (N+2 × 5)
# ===================================================================
def _fill_slide_18(slide, nuit_cam_data, weeks):
    tables = _get_tables(slide)
    if not tables:
        return
    _shape = tables[0]
    t = _shape.table
    nights = sorted(nuit_cam_data.keys())
    needed = 2 + len(nights) + 1  # bandeau + header + data + TOTAL
    _ensure_table_size(t, needed)
    _trim_table_cols(t, 5)  # retire la colonne vide résiduelle à droite
    _set_cell_text(t.cell(0, 0), "Récap par nuit", bold=True, align="center", size=12)
    for ci, h in enumerate(["Nuit", "Date", "Secteur/Rayon", "Allées", "Caméras"]):
        _set_cell_text(t.cell(1, ci), h, bold=True, size=10)
    tot_cam = 0
    tot_allees = set()
    for i, n in enumerate(nights):
        r = i + 2
        d = nuit_cam_data[n]
        _set_cell_text(t.cell(r, 0), f"Nuit {n}", bold=True, size=10)
        _set_cell_text(t.cell(r, 1), _fmt_date(d.get("date")), size=10)
        _set_cell_text(t.cell(r, 2), d.get("sr", ""), align="left", size=9)
        _set_cell_text(t.cell(r, 3), d.get("allees_str", ""), align="left", size=9)
        _set_cell_text(t.cell(r, 4), _num(d.get("cam", 0)), bold=True, size=10)
        color = _color_for_night(n, weeks)
        for ci in range(5):
            _set_cell_fill(t.cell(r, ci), color)
            _force_cell_text_color(t.cell(r, ci), "#000000")
        # Somme totaux
        tot_cam += int(d.get("cam", 0) or 0)
        for a in str(d.get("allees_str") or "").split(","):
            a = a.strip()
            if a:
                tot_allees.add(a)
    # Ligne TOTAL (cf. modèle de référence utilisateur du 10/07/2026)
    tr = 2 + len(nights)
    _set_cell_text(t.cell(tr, 0), "TOTAL", bold=True, align="center", size=10)
    _set_cell_text(t.cell(tr, 1), "", size=10)
    _set_cell_text(t.cell(tr, 2), "", size=10)
    _set_cell_text(t.cell(tr, 3), f"{len(tot_allees)} allées planifiées", bold=True, align="center", size=10)
    _set_cell_text(t.cell(tr, 4), _num(tot_cam), bold=True, align="center", size=10)
    for ci in range(5):
        _set_cell_fill(t.cell(tr, ci), "FEF3C7")  # jaune pâle
        _force_cell_text_color(t.cell(tr, ci), "#000000")
    # Supprime les lignes résiduelles éventuelles + compactage/position (réf.)
    while len(t.rows) > needed:
        last_tr = t._tbl.findall(qn('a:tr'))[-1]
        last_tr.getparent().remove(last_tr)
    _place_table(_shape, 718456, 1403335, 9975669, 280000)


# ===================================================================
# Slide 19 — Détail caméras par allée (split en 2 tables 27×2)
# ===================================================================
def _fill_slide_19(slide, detail_rows, weeks=None):
    """Détail caméras par allée. Rows = list[(nuit, allee, elems_str)].
    Couleur de fond par nuit (cohérence visuelle avec le reste du PPTX).
    Pas de bandeau violet interne — le titre est déjà la title de slide."""
    tables = _get_tables(slide)
    if len(tables) < 2:
        return
    t1_shape, t2_shape = tables[0], tables[1]
    t1, t2 = t1_shape.table, t2_shape.table
    needed = len(detail_rows)
    header_fill = (0xE5, 0xE7, 0xEB)
    WHITE_RGB = (0xFF, 0xFF, 0xFF)

    # Répartition ÉQUILIBRÉE gauche/droite : évite le débordement en bas de
    # slide du tableau gauche et les lignes vides colorées du template à droite.
    if needed <= 14:
        n_left, n_right = needed, 0
    else:
        n_left = (needed + 1) // 2
        n_right = needed - n_left

    def _row_color(n: int) -> str | None:
        if not n or n >= 9999 or not weeks:
            return None
        return _color_for_night(n, weeks)

    def _hex_to_rgb(hx: str):
        hx = hx.lstrip("#")
        return (int(hx[0:2], 16), int(hx[2:4], 16), int(hx[4:6], 16))

    def _trim_rows(tbl, keep: int):
        trs = tbl._tbl.findall(qn('a:tr'))
        while len(trs) > max(1, keep):
            trs[-1].getparent().remove(trs[-1])
            trs.pop()

    # Adapter la hauteur des lignes au nombre d'entrées : plus il y en a,
    # plus les lignes doivent être compactes pour tenir dans la slide (7.5").
    max_side = max(n_left, n_right, 1)
    if max_side <= 15:
        row_h = 220000  # aéré
        cell_size = 8
    elif max_side <= 22:
        row_h = 170000  # standard
        cell_size = 7
    elif max_side <= 30:
        row_h = 145000  # compact
        cell_size = 7
    else:
        row_h = 125000  # ultra-compact (30+ rows)
        cell_size = 6

    def _fill_col(tbl, rows, with_header=True):
        offset = 1 if with_header else 0
        _ensure_table_size(tbl, offset + len(rows))
        _trim_rows(tbl, offset + len(rows))
        if with_header:
            _set_cell_text(tbl.cell(0, 0), "Allées", bold=True, align="left",
                           size=cell_size, color="#000000", fill_rgb=header_fill)
            _set_cell_text(tbl.cell(0, 1), "N° Elements", bold=True, align="left",
                           size=cell_size, color="#000000", fill_rgb=header_fill)
        for i, (n, allee, elems) in enumerate(rows):
            color = _row_color(n)
            fill = _hex_to_rgb(color) if color else WHITE_RGB
            _set_cell_text(tbl.cell(offset + i, 0), allee,
                           bold=False, align="left", size=cell_size, color="#000000",
                           fill_rgb=fill)
            _set_cell_text(tbl.cell(offset + i, 1), elems,
                           align="left", size=max(cell_size - 1, 5), color="#000000",
                           fill_rgb=fill)
            for cc in (tbl.cell(offset + i, 0), tbl.cell(offset + i, 1)):
                for p in cc.text_frame.paragraphs:
                    for r in p.runs:
                        r.font.bold = False
                # Force la couleur noire sur TOUS les runs (le template
                # a des runs gras/rouges hérités que _set_cell_text ne
                # remplace pas fiablement).
                _force_cell_text_color(cc, "#000000")
        # Force également le noir sur les cellules d'en-tête
        if with_header:
            _force_cell_text_color(tbl.cell(0, 0), "#000000")
            _force_cell_text_color(tbl.cell(0, 1), "#000000")
        for tr in tbl._tbl.findall(qn('a:tr')):
            tr.set('h', str(row_h))

    _fill_col(t1, detail_rows[:n_left], with_header=True)
    if n_right == 0:
        # Tableau droit inutile : on le supprime (sinon lignes vides colorées)
        t2_shape._element.getparent().remove(t2_shape._element)
    else:
        _fill_col(t2, detail_rows[n_left:], with_header=True)
        try:
            t2_shape.top = t1_shape.top
        except (TypeError, ValueError):
            pass


# ===================================================================
# Slide 20 — Phasage full consolidé (N+3 × 10)
# ===================================================================
def _fill_slide_20(slide, nuit_es_data, nuit_cam_data, dates_map, weeks):
    tables = _get_tables(slide)
    if not tables:
        return
    t = tables[0].table
    all_n = sorted(set(nuit_es_data.keys()) | set(nuit_cam_data.keys()))
    # 13 colonnes — SA éclatées en 4 (SA 1.5 / SA 2.1 / SA 2.1 frz / SA magasin)
    # cf. modèle de référence utilisateur du 10/07/2026.
    ncols = 13
    _ensure_table_cols(t, ncols, label_cols=1)
    _trim_table_cols(t, ncols)
    # Ratios de colonnes (Allées un peu large, SA compactes, Caméras droite)
    _set_col_widths_by_ratio(t, [8, 6, 6, 5, 5, 5, 6, 12, 5, 7, 12, 8, 6])
    # 3 lignes d'en-tête + N nuits + 1 ligne TOTAL
    needed = 3 + len(all_n) + 1
    _ensure_table_size(t, needed)
    # Ligne 0 : titre principal (fusion sur les 13 colonnes)
    _merge_title_row(t, ncols, row_idx=0)
    _set_cell_text(t.cell(0, 0), "Phasage full — Planning consolidé EEG + Caméras",
                   bold=True, align="center", size=11)
    # Ligne 1 : 3 sous-bandeaux fusionnés
    _merge_cells_range(t, 1, 0, 7)   # « Phasage étiquettes et rails »
    _merge_cells_range(t, 1, 8, 9)   # « Nuit »
    _merge_cells_range(t, 1, 10, 12) # « Phasage caméras »
    _set_cell_text(t.cell(1, 0), "Phasage étiquettes et rails", bold=True, align="center", size=9)
    _set_cell_text(t.cell(1, 8), "Nuit", bold=True, align="center", size=9)
    _set_cell_text(t.cell(1, 10), "Phasage caméras", bold=True, align="center", size=9)
    subs = ["Allées", "ES", "Rails ES", "SA 1.5", "SA 2.1", "SA 2.1 frz", "SA magasin",
            "Secteur/Rayon", "Nuit", "Date",
            "Secteur/Rayon", "Allées", "Caméras"]
    for ci, s in enumerate(subs):
        _set_cell_text(t.cell(2, ci), s, bold=True, size=8)
    # Force noir sur les en-têtes (le template a des runs violets/gris hérités)
    for ci in range(ncols):
        _force_cell_text_color(t.cell(1, ci), "#000000")
        _force_cell_text_color(t.cell(2, ci), "#000000")
    # Data — police compacte pour faire tenir toutes les nuits dans la slide
    tot_es = tot_rails = tot_15 = tot_21 = tot_fz = tot_mag = tot_cam = 0
    for i, n in enumerate(all_n):
        r = i + 3
        es = nuit_es_data.get(n, {})
        cam = nuit_cam_data.get(n, {})
        _set_cell_text(t.cell(r, 0), es.get("allees_str", ""), align="left", size=7)
        _set_cell_text(t.cell(r, 1), _num(es.get("eeg", "")), size=8)
        _set_cell_text(t.cell(r, 2), _num(es.get("rails_es", "")), size=8)
        _set_cell_text(t.cell(r, 3), _num(es.get("sa_inst_15", "") or ""), size=8)
        _set_cell_text(t.cell(r, 4), _num(es.get("sa_inst_21", "") or ""), size=8)
        _set_cell_text(t.cell(r, 5), _num(es.get("sa_inst_freezer", "") or ""), size=8)
        _set_cell_text(t.cell(r, 6), _num(es.get("sa_mag", "") or ""), size=8,
                       italic=True, color="#6B7280")
        _set_cell_text(t.cell(r, 7), es.get("sr", ""), align="left", size=7)
        _set_cell_text(t.cell(r, 8), str(n), bold=True, size=9)
        _set_cell_text(t.cell(r, 9), _fmt_date(dates_map.get(str(n))), size=8)
        _set_cell_text(t.cell(r, 10), cam.get("sr", ""), align="left", size=7)
        _set_cell_text(t.cell(r, 11), cam.get("allees_str", ""), align="left", size=7)
        _set_cell_text(t.cell(r, 12), _num(cam.get("cam", "")), size=8)
        color = _color_for_night(n, weeks)
        for ci in range(ncols):
            _set_cell_fill(t.cell(r, ci), color)
            if ci == 6:  # SA magasin en italique/gris
                _force_cell_text_color(t.cell(r, ci), "#6B7280")
            else:
                _force_cell_text_color(t.cell(r, ci), "#000000")
        # Cumul TOTAL
        def _f(v):
            try: return float(v) if v not in (None, "") else 0
            except (ValueError, TypeError): return 0
        tot_es += _f(es.get("eeg"))
        tot_rails += _f(es.get("rails_es"))
        tot_15 += _f(es.get("sa_inst_15"))
        tot_21 += _f(es.get("sa_inst_21"))
        tot_fz += _f(es.get("sa_inst_freezer"))
        tot_mag += _f(es.get("sa_mag"))
        tot_cam += _f(cam.get("cam"))
    # Ligne TOTAL
    total_row = 3 + len(all_n)
    _set_cell_text(t.cell(total_row, 0), "TOTAL", bold=True, align="left", size=9)
    _set_cell_text(t.cell(total_row, 1), _num(int(tot_es)), bold=True, size=9)
    _set_cell_text(t.cell(total_row, 2), _num(int(tot_rails)), bold=True, size=9)
    _set_cell_text(t.cell(total_row, 3), _num(int(tot_15) or ""), bold=True, size=9)
    _set_cell_text(t.cell(total_row, 4), _num(int(tot_21) or ""), bold=True, size=9)
    _set_cell_text(t.cell(total_row, 5), _num(int(tot_fz) or ""), bold=True, size=9)
    _set_cell_text(t.cell(total_row, 6), _num(int(tot_mag) or ""), bold=True, size=9,
                   italic=True, color="#6B7280")
    _set_cell_text(t.cell(total_row, 7), "", size=9)
    _set_cell_text(t.cell(total_row, 8), f"{len(all_n)} nuits", bold=True, size=9)
    _set_cell_text(t.cell(total_row, 9), "", size=9)
    _set_cell_text(t.cell(total_row, 10), "", size=9)
    _set_cell_text(t.cell(total_row, 11), "", size=9)
    _set_cell_text(t.cell(total_row, 12), _num(int(tot_cam)), bold=True, size=9)
    for ci in range(ncols):
        if ci == 6:
            _force_cell_text_color(t.cell(total_row, ci), "#6B7280")
        else:
            _force_cell_text_color(t.cell(total_row, ci), "#000000")
    # Supprime les rows résiduelles du template
    trs = t._tbl.findall(qn('a:tr'))
    while len(trs) > total_row + 1:
        trs[-1].getparent().remove(trs[-1])
        trs.pop()
    # Force des hauteurs de ligne réduites pour faire tenir tout dans la slide
    for tr in t._tbl.findall(qn('a:tr')):
        tr.set('h', '240000')


# ===================================================================
# Slides "Plan wifi magasin" — feature d'upload retirée mais le template
# contient 2 slides "Plan wifi magasin" (principale + réserve). On garde
# UNE slide vide comme placeholder (ref utilisateur) et on supprime la
# 2e (réserve inutilisée).
# ===================================================================
def _find_wifi_slide_indices(prs) -> list[int]:
    out: list[int] = []
    for i, s in enumerate(prs.slides):
        for sh in s.shapes:
            try:
                if sh.has_text_frame and "plan wifi" in sh.text_frame.text.strip().lower():
                    out.append(i)
                    break
            except Exception:
                continue
    return out


def _remove_wifi_slides(prs) -> None:
    """Supprime UNIQUEMENT la 2e slide 'Plan wifi magasin' (réserve).
    La 1re est conservée comme placeholder vide dans le PPTX (cf. modèle
    de référence utilisateur du 10/07/2026)."""
    idxs = _find_wifi_slide_indices(prs)
    # On conserve la 1re, on supprime les suivantes (en ordre décroissant).
    for idx in sorted(idxs[1:], reverse=True):
        _delete_slide(prs, idx)



# ===================================================================
# (iter48j) Clonage du slide "Accès et logistique" fourni par l'utilisateur
# ===================================================================
def _insert_logistique_slide(prs, after_idx: int) -> bool:
    """Insère le slide unique du fichier `slide_logistique.pptx` dans `prs`
    juste après l'index `after_idx` (0-indexé). Copie shapes + relations
    (images / hyperliens). Retourne True si inséré, False sinon.

    NB : ne re-crée pas les slide layouts — utilise le layout blank du template
    principal, l'apparence peut légèrement différer (bordure du tableau ok).
    """
    if not LOGISTIQUE_SLIDE_PATH.exists():
        logger.warning("Slide logistique introuvable : %s", LOGISTIQUE_SLIDE_PATH)
        return False
    try:
        from copy import deepcopy
        src = Presentation(str(LOGISTIQUE_SLIDE_PATH))
        if not src.slides:
            return False
        src_slide = src.slides[0]
        # Nouveau slide vide dans la destination avec layout blank
        blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
        new_slide = prs.slides.add_slide(blank_layout)
        # Retire les placeholders du blank layout pour éviter les doublons
        for sh in list(new_slide.shapes):
            sh._element.getparent().remove(sh._element)
        # Copie chaque shape XML du slide source (titre, table, textbox…)
        spTree = new_slide.shapes._spTree
        for shape in src_slide.shapes:
            spTree.append(deepcopy(shape._element))
        # Copie les relations (images, hyperliens, tableaux stylés)
        for rel in list(src_slide.part.rels.values()):
            if "notesSlide" in rel.reltype:
                continue
            if rel.is_external:
                new_slide.part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
            else:
                new_slide.part.rels.get_or_add(rel.reltype, rel.target_part)
        # Déplace le slide juste après after_idx (add_slide l'a mis à la fin)
        sldIdLst = prs.slides._sldIdLst
        new_id = sldIdLst[-1]
        sldIdLst.remove(new_id)
        sldIdLst.insert(after_idx + 1, new_id)
        return True
    except Exception as e:
        logger.warning("Impossible d'insérer le slide logistique : %s", e)
        return False


# ===================================================================
# Public entry point
# ===================================================================
def build_pptx(d: dict, *, aggregate_fn, recap_rows: list, summary: dict | None = None,
               detail_cam_rows: list[tuple[int, str, str]] | None = None) -> bytes:
    """Génère le PowerPoint complet à partir des données.

    aggregate_fn(d) → dict avec clés : nuit_es (n -> {date, sr, allees_str, eeg, rails_es, sa, cam}),
    nuit_cam, dates_map, weeks, totals_by_nuit (n -> {eeg, cam, sa}), all_nights, cam_nights.
    """
    if not TEMPLATE_PATH.exists():
        raise FileNotFoundError(f"Template PowerPoint introuvable : {TEMPLATE_PATH}")
    agg = aggregate_fn(d)
    prs = Presentation(str(TEMPLATE_PATH))
    slides = prs.slides
    # NEW TEMPLATE (10/07/2026) : le PPTX fourni par l'utilisateur contient
    # 20 slides (au lieu de 22 précédemment). Les indices des slides à
    # remplir sont décalés de -1 (pas de slide "Accès et logistique" ni de
    # slide wifi de réserve).
    #   idx 7  = Matériel à commander (ex-8)
    #   idx 10 = Plan phasage EEG complet (ex-11)
    #   idx 11 = Récap par nuit (ex-12)
    #   idx 12-15 = Semaines S1-S4 (ex-13-16)
    #   idx 16 = Plan cameras complet (ex-17)
    #   idx 17 = Récap cameras (ex-18)
    #   idx 18 = Détail cameras (ex-19)
    #   idx 19 = Phasage full (ex-20)
    if len(slides) >= 8:
        _fill_slide_8(slides[7], recap_rows)
    # Nb nuits dynamiques pour les titres "(N nuits)"
    nb_nuits_eeg = len(agg.get("all_nights") or [])
    nb_nuits_cam = len(agg.get("cam_nights") or [])
    # Slide 11 (index 10) = Plan phasage EEG complet
    if len(slides) >= 11:
        _replace_nb_nuits_in_title(slides[10], nb_nuits_eeg)
        _fill_slide_11(slides[10], agg["totals_by_nuit"], agg["dates_map"],
                       agg["weeks"], agg["all_nights"])
    # Slide 12 (index 11) = Récap par nuit
    if len(slides) >= 12:
        _replace_nb_nuits_in_title(slides[11], nb_nuits_eeg)
        _fill_slide_12(slides[11], agg["nuit_es"], agg["weeks"],
                       hide_sa_mag=agg.get("hide_sa_mag", False))
    # Slides 13-16 (index 12-15) = semaines S1..S4
    weeks_list = agg["weeks"] or []
    WEEK_SLIDE_INDICES = [12, 13, 14, 15]  # S1, S2, S3, S4
    cursor = 1
    used_week_slides = set()
    for wi, w in enumerate(weeks_list[:4]):
        ww = int(w or 0)
        if ww <= 0:
            continue
        week_nights = list(range(cursor, cursor + ww))
        cursor += ww
        slide_idx = WEEK_SLIDE_INDICES[wi]
        if slide_idx < len(slides):
            _fill_slide_week(slides[slide_idx], wi + 1, week_nights,
                             agg["nuit_es"], agg["totals_by_nuit"],
                             agg["dates_map"], weeks_list,
                             hide_sa_mag=agg.get("hide_sa_mag", False))
            used_week_slides.add(slide_idx)
    # Slide 17 (index 16) = Plan cameras complet
    if len(slides) >= 17 and agg["cam_nights"]:
        _replace_nb_nuits_in_title(slides[16], nb_nuits_cam)
        _fill_slide_17(slides[16], agg["totals_by_nuit"], agg["dates_map"],
                       agg["cam_nights"], weeks_list)
    # Slide 18 (index 17) = Récap caméras par nuit
    if len(slides) >= 18:
        _replace_nb_nuits_in_title(slides[17], nb_nuits_cam)
        _fill_slide_18(slides[17], agg["nuit_cam"], weeks_list)
    # Slide 19 (index 18) = Détail caméras par allée
    if len(slides) >= 19 and detail_cam_rows:
        _fill_slide_19(slides[18], detail_cam_rows, weeks=weeks_list)
    # Slide 20 (index 19) = Phasage full consolidé
    if len(slides) >= 20:
        _fill_slide_20(slides[19], agg["nuit_es"], agg["nuit_cam"],
                       agg["dates_map"], weeks_list)
    # Semaines au-delà de 4 : clone la dernière slide semaine (idx 15)
    if len(weeks_list) > 4 and len(slides) >= 16:
        cursor2 = sum(int(w or 0) for w in weeks_list[:4]) + 1
        for j, w in enumerate(weeks_list[4:]):
            ww = int(w or 0)
            if ww <= 0:
                continue
            week_nights = list(range(cursor2, cursor2 + ww))
            cursor2 += ww
            new_slide = _duplicate_slide(prs, 15, 16 + j)
            _fill_slide_week(new_slide, 5 + j, week_nights,
                             agg["nuit_es"], agg["totals_by_nuit"],
                             agg["dates_map"], weeks_list,
                             hide_sa_mag=agg.get("hide_sa_mag", False))
    # Suppression des slides semaines non utilisées
    unused = sorted([idx for idx in WEEK_SLIDE_INDICES if idx not in used_week_slides], reverse=True)
    for idx in unused:
        if idx < len(slides):
            _delete_slide(prs, idx)
    # Le nouveau template n'a qu'UNE slide "Plan wifi magasin" (placeholder) —
    # on la conserve telle quelle, plus rien à supprimer.
    # Titres remontés + phrase de bas de page réduite (sur toutes les slides).
    try:
        _compact_layout(prs)
    except Exception:
        pass
    # (iter48j) Insertion du slide "Accès et logistique" APRÈS remplissage des
    # autres, pour ne pas décaler les indices utilisés ci-dessus. Position finale
    # attendue : slide 7 (index 6), juste après « Informations Magasin ».
    _insert_logistique_slide(prs, after_idx=5)
    # Sauvegarde en bytes
    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _duplicate_slide(prs, src_index: int, dest_index: int):
    """Clone la slide `src_index` (deep-copy des formes) et la déplace en
    position `dest_index`. Utilisé pour créer une 5e slide "semaine" quand le
    phasage dépasse 4 semaines (jusqu'à 20 nuits = 5 semaines de 4)."""
    src = prs.slides[src_index]
    new_slide = prs.slides.add_slide(src.slide_layout)
    # Retire les placeholders hérités du layout
    for shp in list(new_slide.shapes):
        shp._element.getparent().remove(shp._element)
    # Deep-copy des formes de la slide source (tables + textes ; pas d'images)
    for shp in src.shapes:
        new_slide.shapes._spTree.append(copy.deepcopy(shp._element))
    # Réordonne : déplace la nouvelle slide (en dernier) vers dest_index
    xml_slides = prs.slides._sldIdLst
    ids = list(xml_slides)
    new_id = ids[-1]
    xml_slides.remove(new_id)
    xml_slides.insert(dest_index, new_id)
    return new_slide


def _delete_slide(prs, slide_idx: int):
    """Supprime la slide à l'index donné du PowerPoint en nettoyant aussi
    sa référence et sa relation dans la présentation."""
    xml_slides = prs.slides._sldIdLst
    slides_list = list(xml_slides)
    if slide_idx >= len(slides_list):
        return
    slide_id_elem = slides_list[slide_idx]
    rId = slide_id_elem.get(qn('r:id'))
    # Supprime de la liste sldIdLst
    xml_slides.remove(slide_id_elem)
    # Drop la relation pour que python-pptx oublie la slide
    if rId:
        prs.part.drop_rel(rId)
