"""iter48j — Slide « Accès et logistique » et couleurs des nuits sans doublons.

Bugs signalés par l'utilisateur :
  1. Le slide « Accès et logistique » a disparu du PPTX exporté
  2. Dans les tableaux Récap par nuit, 2 nuits consécutives pouvaient avoir la
     même couleur de fond (bleu-bleu au changement de semaine)

Vérifie :
  A. Le slide « Accès et logistique » est bien inséré en position 7 (après
     « Informations Magasin ») dans le PPTX final
  B. `_color_for_night` cycle strictement bleu → jaune → rose → vert → bleu…
     sur le n° absolu de nuit, indépendamment du découpage par semaine
"""
import sys
from pathlib import Path

sys.path.insert(0, "/app/backend")

from pptx import Presentation
from pptx_export import (
    LOGISTIQUE_SLIDE_PATH, TEMPLATE_PATH,
    WEEK_COLORS_HEX, _color_for_night, _insert_logistique_slide,
)


# ─── A. Couleurs des nuits ───────────────────────────────────────────────────
def test_week_colors_hex_are_all_distinct():
    """Aucune couleur dupliquée dans la palette (le bug initial était que la
    1ère et la 4e couleur étaient identiques → bleu-bleu)."""
    assert len(set(WEEK_COLORS_HEX)) == len(WEEK_COLORS_HEX), WEEK_COLORS_HEX
    assert len(WEEK_COLORS_HEX) == 4


def test_no_two_consecutive_nights_share_color_with_normal_weeks():
    """Semaines normales de 4 nuits : la palette 4 couleurs distinctes garantit
    qu'aucune paire consécutive n'a la même couleur, y compris au changement
    de semaine (vert → bleu)."""
    weeks = [4, 4, 4, 4]
    prev = None
    for n in range(1, 17):
        c = _color_for_night(n, weeks)
        if prev is not None:
            assert c != prev, f"Nuit {n-1} et Nuit {n} même couleur : {c}"
        prev = c


def test_partial_week_restarts_at_blue():
    """Règle métier utilisateur : après une semaine partielle (2 nuits), la
    semaine suivante recommence à bleu. Exemple weeks=[2, 4] :
      N1 bleu, N2 jaune | N3 bleu, N4 jaune, N5 rose, N6 vert
    """
    assert _color_for_night(1, [2, 4]) == "#DBEAFE"  # bleu
    assert _color_for_night(2, [2, 4]) == "#FEF3C7"  # jaune
    assert _color_for_night(3, [2, 4]) == "#DBEAFE"  # bleu (redémarre)
    assert _color_for_night(4, [2, 4]) == "#FEF3C7"  # jaune
    assert _color_for_night(5, [2, 4]) == "#FECACA"  # rose
    assert _color_for_night(6, [2, 4]) == "#DCFCE7"  # vert


def test_color_position_based():
    """Toujours 1ère nuit d'une semaine = bleu (règle utilisateur)."""
    # Semaine unique de 4 nuits : bleu / jaune / rose / vert
    assert _color_for_night(1, [4]) == "#DBEAFE"
    assert _color_for_night(2, [4]) == "#FEF3C7"
    assert _color_for_night(3, [4]) == "#FECACA"
    assert _color_for_night(4, [4]) == "#DCFCE7"


# ─── B. Insertion du slide « Accès et logistique » ───────────────────────────
def test_logistique_slide_asset_present():
    """Le fichier ressource `slide_logistique.pptx` doit exister."""
    assert LOGISTIQUE_SLIDE_PATH.exists(), LOGISTIQUE_SLIDE_PATH


def test_insertion_slide_in_correct_position():
    """Après insertion, le slide « Accès et logistique » doit être en position
    7 (index 6), juste après « Informations Magasin » (index 5)."""
    prs = Presentation(str(TEMPLATE_PATH))
    n_before = len(prs.slides)
    inserted = _insert_logistique_slide(prs, after_idx=5)
    assert inserted is True
    assert len(prs.slides) == n_before + 1
    # Récupère le texte du slide 7 (index 6)
    slide7 = prs.slides[6]
    texts = []
    for sh in slide7.shapes:
        if sh.has_text_frame:
            for para in sh.text_frame.paragraphs:
                for r in para.runs:
                    if r.text.strip():
                        texts.append(r.text)
    joined = " ".join(texts)
    assert "Accès et logistique" in joined, joined


def test_insertion_does_not_corrupt_zip_no_duplicates():
    """(iter48m) Le PPTX final ne doit contenir AUCUN nom de fichier dupliqué
    dans le ZIP — sinon PowerPoint refuse d'ouvrir avec le message
    « PowerPoint a détecté un problème dans le contenu ». Ce test protège
    contre la régression où on clonait aussi slideMasters + slideLayouts."""
    import io, warnings, zipfile
    prs = Presentation(str(TEMPLATE_PATH))
    _insert_logistique_slide(prs, after_idx=5)
    buf = io.BytesIO()
    with warnings.catch_warnings(record=True) as caught:
        warnings.simplefilter("always")
        prs.save(buf)
        dupes = [str(x.message) for x in caught if "Duplicate name" in str(x.message)]
    assert not dupes, f"Noms dupliqués dans le ZIP : {dupes[:5]}"
    # ZIP intègre + réouvrable
    buf.seek(0)
    assert zipfile.ZipFile(buf).testzip() is None
    buf.seek(0)
    prs2 = Presentation(buf)
    assert len(prs2.slides) == 21
