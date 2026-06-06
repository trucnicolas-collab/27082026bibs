"""Tests pour la génération du PowerPoint (CR VT + Plan de phasage).

Vérifient que :
  - L'endpoint /api/dataset/{id}/export-pptx renvoie un fichier valide
  - Les méta-infos magasin sont bien injectées dans les slides clés
  - Les titres "(X nuits)" sont mis à jour dynamiquement
"""
import io
import os
import sys
import pytest

# Permet d'importer pptx_export et server depuis /app/backend
sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".."))

from pptx_export import generate_pptx, _aggregate, _fmt_date_short  # noqa: E402
from pptx import Presentation  # noqa: E402


def _fake_dataset(**over):
    """Construit un dataset minimal pour tester la génération."""
    base = {
        "store_name": "Carrefour Test",
        "store_city": "TestCity",
        "store_code": "T0001",
        "store_address": "1 rue de Test",
        "vt_start_date": "2026-04-27",
        "vt_end_date": "2026-04-29",
        "phasage": {
            "es": {
                "nb_nuits": 4,
                "rows": [
                    {"id": "r1", "allee": "1", "nuit": 1},
                    {"id": "r2", "allee": "2", "nuit": 2},
                    {"id": "r3", "allee": "3", "nuit": 3},
                    {"id": "r4", "allee": "4", "nuit": 4},
                ],
                "weeks": [2, 2],
            },
            "cam": {
                "nb_nuits": 2,
                "start_at_nuit": 3,
                "rows": [
                    {"id": "c1", "allee": "1", "nuit": 1},
                    {"id": "c2", "allee": "2", "nuit": 2},
                ],
            },
            "dates": {
                "1": "2026-05-01", "2": "2026-05-02",
                "3": "2026-05-04", "4": "2026-05-05",
            },
        },
    }
    base.update(over)
    return base


def _fake_summary():
    return {
        "allees": [
            {"uid": "1", "allee": "1", "es_15": 100, "es_21": 50, "rails_es": 10,
             "sa": 0, "sa_15": 0, "sa_21": 0, "cameras": 5, "camera_elems": [1, 2, 3]},
            {"uid": "2", "allee": "2", "es_15": 200, "es_21": 80, "rails_es": 15,
             "sa": 0, "sa_15": 0, "sa_21": 0, "cameras": 8, "camera_elems": [4, 5]},
            {"uid": "3", "allee": "3", "es_15": 150, "es_21": 60, "rails_es": 12,
             "sa": 0, "sa_15": 0, "sa_21": 0, "cameras": 0, "camera_elems": []},
            {"uid": "4", "allee": "4", "es_15": 90, "es_21": 30, "rails_es": 8,
             "sa": 0, "sa_15": 0, "sa_21": 0, "cameras": 0, "camera_elems": []},
        ],
        "totals": {},
    }


def test_generate_pptx_returns_valid_bytes():
    data = generate_pptx(_fake_dataset(), _fake_summary())
    assert isinstance(data, bytes)
    assert len(data) > 1000  # PPT non vide
    prs = Presentation(io.BytesIO(data))
    # Template = 22 slides ; weeks=[2,2] → S3/S4/S5 supprimées → 22 - 3 = 19
    assert len(prs.slides) == 19


def test_store_name_and_code_in_cover():
    data = generate_pptx(_fake_dataset(), _fake_summary())
    prs = Presentation(io.BytesIO(data))
    sub = next((s for s in prs.slides[0].shapes if s.name == "Untertitel 2"), None)
    assert sub is not None
    assert "TestCity" in sub.text_frame.text
    assert "T0001" in sub.text_frame.text


def test_vt_dates_in_slide4():
    data = generate_pptx(_fake_dataset(), _fake_summary())
    prs = Presentation(io.BytesIO(data))
    sub = next((s for s in prs.slides[3].shapes if s.name == "Sous-titre 4"), None)
    assert sub is not None
    assert "27/04/26" in sub.text_frame.text
    assert "29/04/26" in sub.text_frame.text


def test_info_table_filled_slide6():
    data = generate_pptx(_fake_dataset(), _fake_summary())
    prs = Presentation(io.BytesIO(data))
    tbl_shape = next((s for s in prs.slides[5].shapes
                      if s.has_table and s.name == "Tableau 5"), None)
    assert tbl_shape is not None
    cells_text = [
        (row.cells[0].text.strip().lower(), row.cells[1].text.strip())
        for row in tbl_shape.table.rows
    ]
    mapping = dict(cells_text)
    assert mapping["nom magasin"] == "Carrefour Test"
    assert mapping["code magasin"] == "T0001"
    assert "1 rue de Test" in mapping["adresse"]


def test_dynamic_night_count_in_titles():
    """(X nuits) doit refléter le nb_nuits actuel du dataset."""
    data = generate_pptx(_fake_dataset(), _fake_summary())
    prs = Presentation(io.BytesIO(data))
    # Slide 12 (idx 11) : "Plan de phasage EEG et rails complet par nuit (4 nuits)"
    title11 = next((s for s in prs.slides[11].shapes if s.name == "Titre 1"), None)
    assert title11 is not None
    assert "(4 nuits)" in title11.text_frame.text
    # Avec weeks=[2,2], S3-S5 sont supprimées donc le slide "cam complet"
    # n'est plus à l'index 17 mais à l'index 17 - 3 = 14.
    # On le retrouve par recherche du shape ZoneTexte 38 contenant "(X nuits)".
    found = False
    for s in prs.slides:
        for sh in s.shapes:
            if sh.name == "ZoneTexte 38" and sh.has_text_frame and "(2 nuits)" in sh.text_frame.text:
                found = True
                break
        if found:
            break
    assert found, "Le titre '(2 nuits)' caméras doit être présent dans le PPT"


def test_week_date_table_filled():
    """S2 doit afficher les dates des nuits 3-4 (semaine 2).
    Note : avec weeks=[2,2], la slide S2 est à l'index 14 (15ème slide).
    """
    data = generate_pptx(_fake_dataset(), _fake_summary())
    prs = Presentation(io.BytesIO(data))
    # On recherche par nom de table (Tableau 2 = S2)
    tbl_shape = None
    for s in prs.slides:
        sh = next((x for x in s.shapes
                   if x.has_table and x.name == "Tableau 2"), None)
        if sh:
            tbl_shape = sh
            break
    assert tbl_shape is not None
    cells = [c.text.strip() for c in tbl_shape.table.rows[0].cells]
    assert cells[0] == "Nuit 3"
    assert cells[1] == "Nuit 4"
    dates_row = [c.text.strip() for c in tbl_shape.table.rows[1].cells]
    assert dates_row[0] == "04/05/2026"
    assert dates_row[1] == "05/05/2026"


def test_unused_week_slides_are_removed():
    """Si seules 3 semaines sont définies, les slides S4 et S5 doivent
    être physiquement supprimées du PPT."""
    d = _fake_dataset()
    d["phasage"]["es"]["nb_nuits"] = 6
    d["phasage"]["es"]["weeks"] = [2, 2, 2]  # 3 semaines seulement
    d["phasage"]["es"]["rows"] = [
        {"id": f"r{i}", "allee": str(i), "nuit": i} for i in range(1, 7)
    ]
    summary = _fake_summary()
    summary["allees"].extend([
        {"uid": "5", "allee": "5", "es_15": 50, "es_21": 25, "rails_es": 5,
         "sa": 0, "sa_15": 0, "sa_21": 0, "cameras": 0, "camera_elems": []},
        {"uid": "6", "allee": "6", "es_15": 80, "es_21": 40, "rails_es": 7,
         "sa": 0, "sa_15": 0, "sa_21": 0, "cameras": 0, "camera_elems": []},
    ])
    data = generate_pptx(d, summary)
    prs = Presentation(io.BytesIO(data))
    # 22 (template) - 2 (S4, S5 supprimées) = 20
    assert len(prs.slides) == 20
    # On ne doit plus trouver de slide titrée "– S4" ou "– S5"
    all_text = []
    for s in prs.slides:
        for sh in s.shapes:
            if sh.has_text_frame:
                all_text.append(sh.text_frame.text)
    joined = " ".join(all_text)
    assert "– S4" not in joined
    assert "– S5" not in joined
    # Par contre S1/S2/S3 doivent être présentes
    assert "– S3" in joined


def test_cam_recap_table_filled():
    """Le tableau Récap caméras doit contenir les nuits cam (3-4)."""
    data = generate_pptx(_fake_dataset(), _fake_summary())
    prs = Presentation(io.BytesIO(data))
    cam_tbl = None
    for s in prs.slides:
        sh = next((x for x in s.shapes
                   if x.has_table and x.name == "Tableau 1"), None)
        if sh:
            cam_tbl = sh
            break
    assert cam_tbl is not None
    assert cam_tbl.table.rows[2].cells[0].text.strip() == "Nuit 3"
    assert cam_tbl.table.rows[3].cells[0].text.strip() == "Nuit 4"
    total_row = cam_tbl.table.rows[len(cam_tbl.table.rows) - 1]
    assert "TOTAL" in total_row.cells[0].text.upper()
    assert total_row.cells[2].text.strip() == "13"


def test_global_recap_table_filled():
    """Grand tableau récap doit avoir des données ES et cam."""
    data = generate_pptx(_fake_dataset(), _fake_summary())
    prs = Presentation(io.BytesIO(data))
    tbl = None
    for s in prs.slides:
        sh = next((x for x in s.shapes
                   if x.has_table and x.name == "Tableau 4"), None)
        if sh:
            tbl = sh
            break
    assert tbl is not None
    row1 = tbl.table.rows[2]
    assert row1.cells[4].text.strip() == "1"
    assert row1.cells[0].text.strip() == "1"
    total_row = tbl.table.rows[len(tbl.table.rows) - 1]
    assert "760" in total_row.cells[1].text


def test_fallback_when_store_info_missing():
    """Le PPT doit se générer même sans aucune méta-info magasin."""
    d = {"phasage": {"es": {"nb_nuits": 0, "rows": []},
                     "cam": {"nb_nuits": 0, "rows": [], "start_at_nuit": 5},
                     "dates": {}}}
    data = generate_pptx(d, {"allees": [], "totals": {}})
    prs = Presentation(io.BytesIO(data))
    sub = next((s for s in prs.slides[0].shapes if s.name == "Untertitel 2"), None)
    # Doit être vide / espace, pas planté
    assert sub.text_frame.text.strip() == ""


def test_fmt_date_short():
    assert _fmt_date_short("2026-04-27") == "27/04/26"
    assert _fmt_date_short("") == ""
    assert _fmt_date_short(None) == ""
    assert _fmt_date_short("not-a-date") == ""


def test_aggregate_basic():
    agg = _aggregate(_fake_dataset(), _fake_summary())
    assert agg["nb_nuits_es"] == 4
    assert agg["nb_nuits_cam"] == 2
    assert agg["cam_start_at"] == 3
    # 1ère nuit ES doit contenir l'allée "1" et avoir es=150 (100+50)
    assert agg["es_per_nuit"][1]["es"] == 150
    # Nuit cam globale 3 (start_at + 1 - 1 = 3) contient allée "1" avec 5 cam
    assert agg["cam_per_nuit"][3]["cam"] == 5
