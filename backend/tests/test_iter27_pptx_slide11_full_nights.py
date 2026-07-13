"""Iter27 — Fix PPTX slide 11 : toutes les nuits du phasage sont exportées.

Bug : le slide 11 (Plan de phasage EEG et rails par nuit complet) n'affichait
que les nuits ayant au moins une allée assignée dans le phasage ES. Résultat :
les dernières nuits (ex : 17, 18) apparaissaient VIDES sur la droite du
tableau alors que le titre annonçait « 18 nuits » et que la config phasage
avait bien nb_nuits=18.

Fix : `all_nights = sorted(all keys) | range(1, nb_nuits+1)` dans l'adapter,
+ nettoyage des attributs de fusion (gridSpan/hMerge) hérités du template
lors du clonage de colonnes.
"""
import io
import os
import requests
from pptx import Presentation
from pptx.oxml.ns import qn

BASE_URL = os.environ.get("REACT_APP_BACKEND_URL", "").rstrip("/")
API = f"{BASE_URL}/api"

ADMIN_EMAIL = "admin@vusion.local"
ADMIN_PWD = "admin123"
DATASET_ID = "fd15443f-6d2a-4cef-bd72-c56bb29e9c42"


def _session():
    s = requests.Session()
    r = s.post(f"{API}/auth/login", json={"email": ADMIN_EMAIL, "password": ADMIN_PWD})
    assert r.status_code == 200, r.text
    return s


def _export_pptx():
    s = _session()
    r = s.get(f"{API}/export-pptx/{DATASET_ID}")
    assert r.status_code == 200, r.text
    return Presentation(io.BytesIO(r.content))


def _slide11_table(prs):
    slide = prs.slides[10]
    for shape in slide.shapes:
        if shape.has_table:
            return shape.table
    raise AssertionError("Slide 11 sans tableau")


def test_slide11_covers_all_phasage_nights():
    """Le slide 11 doit contenir 1 + nb_nuits colonnes (1 label + N nuits)."""
    prs = _export_pptx()
    t = _slide11_table(prs)
    # Dataset test : nb_es = 10
    assert len(t.columns) >= 11, f"Slide 11 devrait avoir >= 11 cols, obtenu {len(t.columns)}"
    # La ligne d'en-tête doit lister Nuit 1..Nuit 10 (au moins)
    header_row = t._tbl.findall(qn('a:tr'))[0]
    tcs = header_row.findall(qn('a:tc'))
    texts = []
    for tc in tcs:
        el = tc.find('.//' + qn('a:t'))
        texts.append(el.text if el is not None else "")
    # Skip 1er label
    night_labels = [x for x in texts[1:] if x and x.startswith("Nuit ")]
    numbers = sorted(int(x.split()[1]) for x in night_labels)
    assert 1 in numbers and 10 in numbers, f"Nuit 1 et 10 attendues, obtenu {numbers}"


def test_slide11_cells_are_not_merged_ghost():
    """Les nouvelles cellules clonées ne doivent pas avoir d'attribut de
    fusion (gridSpan/hMerge/rowSpan/vMerge) hérité du template."""
    prs = _export_pptx()
    t = _slide11_table(prs)
    for row_idx, tr in enumerate(t._tbl.findall(qn('a:tr'))):
        tcs = tr.findall(qn('a:tc'))
        for col_idx, tc in enumerate(tcs):
            for attr in ('gridSpan', 'hMerge', 'rowSpan', 'vMerge'):
                v = tc.get(attr)
                assert not v, \
                    f"Cellule row={row_idx} col={col_idx} a {attr}={v} — fantôme probable"


def test_slide11_eeg_row_has_values_for_all_nights():
    """La ligne EEG ES+SA doit contenir des cellules texte pour chaque nuit
    (même si la valeur est 0 pour les nuits sans allée assignée)."""
    prs = _export_pptx()
    t = _slide11_table(prs)
    tr = t._tbl.findall(qn('a:tr'))[2]  # Row 2 = EEG ES+SA
    tcs = tr.findall(qn('a:tc'))
    label = tcs[0].find('.//' + qn('a:t'))
    assert label is not None and label.text == "EEG ES+SA"
    # Toutes les cellules de données doivent avoir un run text (même si "0")
    for i, tc in enumerate(tcs[1:], start=1):
        el = tc.find('.//' + qn('a:t'))
        assert el is not None, f"Colonne {i} n'a pas de <a:t>"
        assert el.text is not None, f"Colonne {i} texte None"
